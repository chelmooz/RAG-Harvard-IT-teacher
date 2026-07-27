"""
Document Processor v5.8 ALL-IN-ONE — AMD BC-250 (Cyan Skillfish)
=================================================================
Formats supportés et organisation des uploads :
  pdf/    → .pdf
  video/  → .mp4
  audio/  → .mp3, .wav
  office/ → .docx, .pptx, .xlsx
  text/   → .txt, .md
  autre/  → tout autre format accepté

Les sous-dossiers sont créés automatiquement au démarrage.
Chaque fichier uploadé est rangé dans le bon sous-dossier selon son extension.

Optimisations BC-250 :
  - Python 3.13 asyncio.TaskGroup : parallélisme natif
  - Chunking : RecursiveCharacterTextSplitter (0 GPU, 0 modèle CANINE)
  - Batch upload mémoire : lecture en chunks de 8 Mo (GDDR6 unifiée)
  - asyncio.to_thread pour toutes les extractions CPU-bound
  - Whisper singleton : chargé une seule fois, réutilisé pour tous les audios/vidéos
"""

import asyncio
import os
import uuid
from pathlib import Path
from typing import List, Dict, Any, Optional
from loguru import logger

try:
    from pypdf import PdfReader
except ImportError:
    import PyPDF2 as _pypdf2
    PdfReader = _pypdf2.PdfReader

from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl

try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
except ImportError:
    from langchain.text_splitter import RecursiveCharacterTextSplitter


class DocumentProcessor:
    """Processeur de documents multi-formats optimisé Python 3.13 + AMD BC-250."""

    _GPU_SEM: asyncio.Semaphore = asyncio.Semaphore(2)

    SUPPORTED_FORMATS = {
        ".pdf":  "pdf",
        ".txt":  "text",
        ".md":   "text",
        ".docx": "docx",
        ".pptx": "pptx",
        ".xlsx": "xlsx",
        ".mp3":  "audio",
        ".mp4":  "video",
        ".wav":  "audio",
    }

    # Correspondance type → sous-dossier upload
    # Les fichiers sont rangés automatiquement selon leur extension.
    UPLOAD_SUBDIR = {
        "pdf":   "pdf",
        "text":  "text",
        "docx":  "office",
        "pptx":  "office",
        "xlsx":  "office",
        "audio": "audio",
        "video": "video",
    }

    _SPLITTER = RecursiveCharacterTextSplitter(
        chunk_size=400,
        chunk_overlap=80,
        separators=["\n\n", "\n", ". ", " ", ""],
        length_function=len,
    )

    def __init__(self, upload_dir: str = "/app/data/uploads"):
        self.upload_dir = Path(upload_dir)
        self.upload_dir.mkdir(parents=True, exist_ok=True)

        # Créer les sous-dossiers automatiquement au démarrage
        for subdir in ("pdf", "video", "audio", "office", "text", "autre"):
            (self.upload_dir / subdir).mkdir(exist_ok=True)
        logger.info(
            f"📁 Dossiers uploads prêts : "
            f"pdf/ video/ audio/ office/ text/ autre/"
        )

        self._whisper_model = None
        self._whisper_device: Optional[str] = None

    async def save_file(self, file, file_id: str) -> str:
        ext = Path(file.filename).suffix.lower()
        if ext not in self.SUPPORTED_FORMATS:
            raise ValueError(f"Extension non autorisée : {ext}")

        # Déterminer le sous-dossier selon le type de fichier
        file_type = self.SUPPORTED_FORMATS[ext]
        subdir    = self.UPLOAD_SUBDIR.get(file_type, "autre")
        dest_dir  = self.upload_dir / subdir
        dest_dir.mkdir(exist_ok=True)   # sécurité si créé après démarrage

        file_path = dest_dir / f"{file_id}{ext}"
        CHUNK = 8 * 1024 * 1024  # 8 Mo

        with open(file_path, "wb") as f:
            while True:
                data = await file.read(CHUNK)
                if not data:
                    break
                f.write(data)

        size = file_path.stat().st_size
        logger.info(
            f"💾 Sauvegardé : uploads/{subdir}/{file_path.name} "
            f"({size // 1024} Ko)"
        )
        return str(file_path)

    async def process_document(
        self, file_path: str, filename: str
    ) -> List[Dict[str, Any]]:
        ext = Path(file_path).suffix.lower()
        if ext not in self.SUPPORTED_FORMATS:
            raise ValueError(f"Format non supporté : {ext}")

        file_type = self.SUPPORTED_FORMATS[ext]
        extractor = {
            "pdf":   self._extract_pdf,
            "text":  self._extract_text,
            "docx":  self._extract_docx,
            "pptx":  self._extract_pptx,
            "xlsx":  self._extract_xlsx,
            "audio": self._extract_audio,
            "video": self._extract_audio,
        }[file_type]

        text = await asyncio.to_thread(extractor, file_path)
        chunks = self._chunk_text(text, filename, file_type)
        logger.info(f"✅ {len(chunks)} chunks créés pour « {filename} »")
        return chunks

    def _extract_pdf(self, file_path: str) -> str:
        text_parts = []
        try:
            with open(file_path, "rb") as f:
                reader = PdfReader(f)
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        text_parts.append(t)
        except Exception as e:
            logger.error(f"❌ Extraction PDF : {e}")
        return "\n\n".join(text_parts)

    def _extract_text(self, file_path: str) -> str:
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
        except Exception as e:
            logger.error(f"❌ Extraction texte : {e}")
            return ""

    def _extract_docx(self, file_path: str) -> str:
        try:
            doc = DocxDocument(file_path)
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            logger.error(f"❌ Extraction DOCX : {e}")
            return ""

    def _extract_pptx(self, file_path: str) -> str:
        try:
            prs = Presentation(file_path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
            return "\n\n".join(parts)
        except Exception as e:
            logger.error(f"❌ Extraction PPTX : {e}")
            return ""

    def _extract_xlsx(self, file_path: str) -> str:
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
            rows = []
            for sheet in wb:
                for row in sheet.iter_rows(values_only=True):
                    r = " | ".join(str(c) for c in row if c is not None)
                    if r:
                        rows.append(r)
            wb.close()
            return "\n".join(rows)
        except Exception as e:
            logger.error(f"❌ Extraction XLSX : {e}")
            return ""

    def _get_whisper_model(self):
        import whisper
        import torch as _torch

        current_device = "cuda" if _torch.cuda.is_available() else "cpu"

        if self._whisper_model is None or self._whisper_device != current_device:
            if self._whisper_model is not None:
                logger.info("🔄 Whisper : rechargement (changement de device)")
            logger.info(f"🎤 Chargement Whisper base ({current_device})...")
            self._whisper_model = whisper.load_model("base", device=current_device)
            self._whisper_device = current_device
            logger.info(f"✅ Whisper base chargé sur {current_device}")

        return self._whisper_model

    def _extract_audio(self, file_path: str) -> str:
        try:
            model = self._get_whisper_model()
            fp16 = (self._whisper_device == "cuda")
            logger.info(f"🎙️  Transcription {file_path} (fp16={fp16})...")
            result = model.transcribe(file_path, language="fr", fp16=fp16)
            text = result.get("text", "")
            logger.info(f"✅ Transcription terminée ({len(text)} chars)")
            return text
        except ImportError:
            logger.warning("⚠️  Whisper non installé — transcription ignorée")
            return "[Transcription non disponible — installer openai-whisper]"
        except Exception as e:
            logger.error(f"❌ Transcription : {e}")
            return ""

    def unload_whisper(self):
        if self._whisper_model is not None:
            import torch as _torch
            del self._whisper_model
            self._whisper_model = None
            self._whisper_device = None
            if _torch.cuda.is_available():
                _torch.cuda.empty_cache()
            logger.info("🗑️  Whisper déchargé (VRAM libérée)")

    def _chunk_text(
        self,
        text: str,
        source: str,
        file_type: str,
    ) -> List[Dict[str, Any]]:
        if not text.strip():
            return []

        raw_chunks = self._SPLITTER.split_text(text)

        return [
            {
                "text": chunk,
                "metadata": {
                    "source": source,
                    "file_type": file_type,
                    "chunking_method": "recursive_char",
                    "chunk_size": len(chunk),
                },
            }
            for chunk in raw_chunks
            if chunk.strip()
        ]

    async def index_directory(
        self,
        directory: str,
        rag_engine,
    ) -> Dict[str, Any]:
        dir_path = Path(directory)
        if not dir_path.exists():
            raise ValueError(f"Répertoire inexistant : {directory}")

        files = [
            p for p in dir_path.rglob("*")
            if p.is_file() and p.suffix.lower() in self.SUPPORTED_FORMATS
        ]

        stats = {
            "total_files": len(files),
            "processed": 0,
            "failed": 0,
            "total_chunks": 0,
        }

        async def _process_one(file_path: Path):
            async with DocumentProcessor._GPU_SEM:
                logger.info(f"🔒 GPU Sémaphore acquis pour {file_path.name}")
                try:
                    chunks = await self.process_document(str(file_path), file_path.name)
                    await rag_engine.index_chunks(
                        chunks, str(file_path.stem), file_path.name
                    )
                    stats["processed"] += 1
                    stats["total_chunks"] += len(chunks)
                    logger.info(f"✅ {file_path.name} : {len(chunks)} chunks indexés")
                except Exception as e:
                    stats["failed"] += 1
                    logger.error(f"❌ {file_path.name} : {e}")
                finally:
                    logger.info(f"🔓 GPU Sémaphore libéré pour {file_path.name}")

        async with asyncio.TaskGroup() as tg:
            for fp in files:
                tg.create_task(_process_one(fp))

        logger.info(
            f"📊 Indexation terminée : {stats['processed']}/{stats['total_files']} "
            f"fichiers | {stats['total_chunks']} chunks | {stats['failed']} erreurs"
        )
        return stats
