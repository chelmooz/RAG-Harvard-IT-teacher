"""
Document Processor v5.1 — AMD BC-250 (Cyan Skillfish)
=======================================================
CORRECTIFS v5.1 appliqués :
  - FIX W8  : Whisper device via torch.cuda.is_available() (plus robuste que HSA env var)
  - FIX W9  : Singleton Whisper — self._whisper_model évite le rechargement à chaque fichier
  - FIX W10 : Import PyPDF2 remplacé par pypdf (pypdf est le successeur maintenu de PyPDF2)
  - FIX W11 : RecursiveCharacterTextSplitter depuis langchain_text_splitters (package dédié)

OPTIMISATIONS CONSERVÉES :
  1. Python 3.13 asyncio.TaskGroup : parallélisme natif sans ThreadPoolExecutor
  2. Chunking : RecursiveCharacterTextSplitter remplace WTPSplit (0 GPU, 0 modèle CANINE)
  3. Batch d'upload mémoire : lecture en chunks de 8 Mo (pression mémoire GDDR6 maîtrisée)
  4. asyncio.to_thread pour toutes les extractions CPU-bound
"""

import asyncio
import os
import uuid
from pathlib import Path
from typing import List, Dict, Any, Optional
from loguru import logger

# ── Extraction documents ───────────────────────────────────────────────────────
# pypdf est le successeur officiel de PyPDF2 (même API, mieux maintenu)
try:
    from pypdf import PdfReader
except ImportError:
    import PyPDF2 as _pypdf2
    PdfReader = _pypdf2.PdfReader  # fallback si pypdf non installé

from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl

# ── Chunking ───────────────────────────────────────────────────────────────────
# Essai langchain_text_splitters (package dédié) puis fallback langchain complet
try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
except ImportError:
    from langchain.text_splitter import RecursiveCharacterTextSplitter


class DocumentProcessor:
    """Processeur de documents multi-formats optimisé Python 3.13 + AMD BC-250."""

    SUPPORTED_FORMATS = {
        ".pdf":  "pdf",
        ".txt":  "text",
        ".md":   "text",
        ".docx": "docx",
        ".pptx": "pptx",
        ".xlsx": "xlsx",
        ".mp3":  "audio",
        ".mp4":  "video",
        ".wav":  "audio",
    }

    # Séparateurs hiérarchiques : paragraphe > ligne > phrase > mot > char
    # chunk_size=400 : embeddings plus précis pour le recall pgvector HNSW
    # chunk_overlap=80 : préserve le contexte inter-chunks
    _SPLITTER = RecursiveCharacterTextSplitter(
        chunk_size=400,
        chunk_overlap=80,
        separators=["\n\n", "\n", ". ", " ", ""],
        length_function=len,
    )

    def __init__(self, upload_dir: str = "/app/data/uploads"):
        self.upload_dir = Path(upload_dir)
        self.upload_dir.mkdir(parents=True, exist_ok=True)

        # FIX W9 : singleton Whisper — chargé une seule fois, conservé en mémoire
        # Évite ~3s de rechargement + fluctuations VRAM à chaque fichier audio
        self._whisper_model = None
        self._whisper_device: Optional[str] = None

    # ── Sauvegarde fichier ─────────────────────────────────────────────────────

    async def save_file(self, file, file_id: str) -> str:
        """
        Sauvegarde un fichier uploadé en streaming par chunks de 8 Mo.

        POURQUOI 8 Mo ?
        La GDDR6 unifiée BC-250 traite les accès mémoire en bursts de 256 bits.
        Lire le fichier entier d'un coup peut saturer les 16 Go si plusieurs
        uploads sont simultanés. 8 Mo = bon compromis débit / pression mémoire.
        """
        ext = Path(file.filename).suffix.lower()
        if ext not in self.SUPPORTED_FORMATS:
            raise ValueError(f"Extension non autorisée : {ext}")

        # file_id est toujours un UUID (généré par l'appelant) → pas de path traversal
        file_path = self.upload_dir / f"{file_id}{ext}"
        CHUNK = 8 * 1024 * 1024  # 8 Mo

        with open(file_path, "wb") as f:
            while True:
                data = await file.read(CHUNK)
                if not data:
                    break
                f.write(data)

        size = file_path.stat().st_size
        logger.info(f"💾 Sauvegardé : {file_path.name} ({size // 1024} Ko)")
        return str(file_path)

    # ── Pipeline principal ─────────────────────────────────────────────────────

    async def process_document(
        self, file_path: str, filename: str
    ) -> List[Dict[str, Any]]:
        """
        Pipeline complet : extraction + chunking.

        Les extractions lourdes (PDF, DOCX…) sont CPU-bound.
        asyncio.to_thread() les délègue à un thread OS pour ne pas bloquer
        la boucle d'événements FastAPI pendant le traitement.
        """
        ext = Path(file_path).suffix.lower()
        if ext not in self.SUPPORTED_FORMATS:
            raise ValueError(f"Format non supporté : {ext}")

        file_type = self.SUPPORTED_FORMATS[ext]
        extractor = {
            "pdf":   self._extract_pdf,
            "text":  self._extract_text,
            "docx":  self._extract_docx,
            "pptx":  self._extract_pptx,
            "xlsx":  self._extract_xlsx,
            "audio": self._extract_audio,
            "video": self._extract_audio,
        }[file_type]

        text = await asyncio.to_thread(extractor, file_path)
        chunks = self._chunk_text(text, filename, file_type)
        logger.info(f"✅ {len(chunks)} chunks créés pour « {filename} »")
        return chunks

    # ── Extracteurs (synchrones, appelés via asyncio.to_thread) ───────────────

    def _extract_pdf(self, file_path: str) -> str:
        """
        Extraction PDF avec pypdf (successeur maintenu de PyPDF2).
        Pour des PDF scannés (sans couche texte), ajouter pytesseract.
        """
        text_parts = []
        try:
            with open(file_path, "rb") as f:
                reader = PdfReader(f)
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        text_parts.append(t)
        except Exception as e:
            logger.error(f"❌ Extraction PDF : {e}")
        return "\n\n".join(text_parts)

    def _extract_text(self, file_path: str) -> str:
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
        except Exception as e:
            logger.error(f"❌ Extraction texte : {e}")
            return ""

    def _extract_docx(self, file_path: str) -> str:
        try:
            doc = DocxDocument(file_path)
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            logger.error(f"❌ Extraction DOCX : {e}")
            return ""

    def _extract_pptx(self, file_path: str) -> str:
        try:
            prs = Presentation(file_path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
            return "\n\n".join(parts)
        except Exception as e:
            logger.error(f"❌ Extraction PPTX : {e}")
            return ""

    def _extract_xlsx(self, file_path: str) -> str:
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
            rows = []
            for sheet in wb:
                for row in sheet.iter_rows(values_only=True):
                    r = " | ".join(str(c) for c in row if c is not None)
                    if r:
                        rows.append(r)
            wb.close()
            return "\n".join(rows)
        except Exception as e:
            logger.error(f"❌ Extraction XLSX : {e}")
            return ""

    def _get_whisper_model(self):
        """
        FIX W9 : Singleton Whisper.
        Charge le modèle une seule fois et le conserve dans self._whisper_model.
        Évite ~3s de chargement + pic VRAM à chaque fichier audio.
        FIX W8 : device via torch.cuda.is_available() (robuste) et non
        la variable HSA_OVERRIDE_GFX_VERSION qui peut être définie sans GPU réel.
        """
        import whisper
        import torch as _torch

        # FIX W8 : même logique que _get_device() dans rag_engine.py
        current_device = "cuda" if _torch.cuda.is_available() else "cpu"

        # Recharger si le device a changé (rare, mais possible entre tests)
        if self._whisper_model is None or self._whisper_device != current_device:
            if self._whisper_model is not None:
                logger.info("🔄 Whisper : rechargement (changement de device)")
            logger.info(f"🎤 Chargement Whisper base ({current_device})...")
            self._whisper_model = whisper.load_model("base", device=current_device)
            self._whisper_device = current_device
            logger.info(f"✅ Whisper base chargé sur {current_device}")

        return self._whisper_model

    def _extract_audio(self, file_path: str) -> str:
        """
        Transcription audio/vidéo via Whisper base (145 Mo).

        OPTIMISATION BC-250 :
        - Modèle "base" : tient dans 12 Go VRAM avec Mistral Q4 chargé.
        - fp16 uniquement si CUDA disponible (évite les erreurs CPU fp16).
        - FIX W9 : singleton — le modèle est conservé entre les appels.
        """
        try:
            model = self._get_whisper_model()
            fp16 = (self._whisper_device == "cuda")
            logger.info(f"🎙️  Transcription {file_path} (fp16={fp16})...")
            result = model.transcribe(file_path, language="fr", fp16=fp16)
            text = result.get("text", "")
            logger.info(f"✅ Transcription terminée ({len(text)} chars)")
            return text
        except ImportError:
            logger.warning("⚠️  Whisper non installé — transcription ignorée")
            return "[Transcription non disponible — installer openai-whisper]"
        except Exception as e:
            logger.error(f"❌ Transcription : {e}")
            return ""

    def unload_whisper(self):
        """
        Libère explicitement le modèle Whisper de la VRAM.
        À appeler avant le fine-tuning ou si la VRAM est sous pression.
        """
        if self._whisper_model is not None:
            import torch as _torch
            del self._whisper_model
            self._whisper_model = None
            self._whisper_device = None
            if _torch.cuda.is_available():
                _torch.cuda.empty_cache()
            logger.info("🗑️  Whisper déchargé (VRAM libérée)")

    # ── Chunking ───────────────────────────────────────────────────────────────

    def _chunk_text(
        self,
        text: str,
        source: str,
        file_type: str,
    ) -> List[Dict[str, Any]]:
        """
        Découpe le texte avec RecursiveCharacterTextSplitter.

        Séparateurs hiérarchiques [\\n\\n, \\n, '. ', ' ', ''] :
        respect de la structure du document (paragraphes > lignes > phrases).
        Zéro GPU requis — libère les 24 CUs RDNA2 pour les embeddings.
        """
        if not text.strip():
            return []

        raw_chunks = self._SPLITTER.split_text(text)

        return [
            {
                "text": chunk,
                "metadata": {
                    "source": source,
                    "file_type": file_type,
                    "chunking_method": "recursive_char",
                    "chunk_size": len(chunk),
                },
            }
            for chunk in raw_chunks
            if chunk.strip()
        ]

    # ── Indexation répertoire (Python 3.13 TaskGroup) ─────────────────────────

    async def index_directory(
        self,
        directory: str,
        rag_engine,
    ) -> Dict[str, Any]:
        """
        Indexe tous les fichiers d'un répertoire EN PARALLÈLE.

        asyncio.TaskGroup (Python 3.13) lance toutes les tâches simultanément.
        Chaque tâche = extraction (thread) + encodage GPU batch + INSERT asyncpg.
        Sur BC-250 : 6 cœurs Zen 2 extraient les PDF pendant que
        les 24 CUs RDNA2 encodent les embeddings du fichier précédent.

        try/except dans _process_one : un fichier défaillant n'annule pas tout.
        """
        dir_path = Path(directory).resolve(strict=False)
        # Whitelist : seuls les sous-répertoires de upload_dir sont autorisés
        allowed = self.upload_dir.resolve()
        if not str(dir_path).startswith(str(allowed)):
            raise ValueError(
                f"Chemin refusé : {dir_path} — seuls les sous-répertoires de "
                f"{allowed} sont autorisés"
            )
        if not dir_path.exists():
            raise ValueError(f"Répertoire inexistant : {directory}")

        files = [
            p for p in dir_path.rglob("*")
            if p.is_file() and p.suffix.lower() in self.SUPPORTED_FORMATS
        ]

        stats = {
            "total_files": len(files),
            "processed": 0,
            "failed": 0,
            "total_chunks": 0,
        }

        async def _process_one(file_path: Path):
            try:
                chunks = await self.process_document(str(file_path), file_path.name)
                await rag_engine.index_chunks(
                    chunks, str(file_path.stem), file_path.name
                )
                stats["processed"] += 1
                stats["total_chunks"] += len(chunks)
                logger.info(f"✅ {file_path.name} : {len(chunks)} chunks indexés")
            except Exception as e:
                stats["failed"] += 1
                logger.error(f"❌ {file_path.name} : {e}")

        # TaskGroup Python 3.13 : annulation propre si exception non capturée
        async with asyncio.TaskGroup() as tg:
            for fp in files:
                tg.create_task(_process_one(fp))

        logger.info(
            f"📊 Indexation terminée : {stats['processed']}/{stats['total_files']} "
            f"fichiers | {stats['total_chunks']} chunks | {stats['failed']} erreurs"
        )
        return stats
