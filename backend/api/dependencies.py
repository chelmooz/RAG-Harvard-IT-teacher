"""
Dependency Injection Container for Prof IA v6.1.
Wires concrete implementations to protocols.
"""
from functools import lru_cache

from fastapi import Depends

from .config import get_settings
from .database import close_db, get_db
from .document_processor import DocumentProcessor
from .protocols import (
    Chunker,
    DocumentExtractor,
    EmbeddingProvider,
    LLMClient,
    Transcriber,
    VectorStore,
)
from .rag_engine import Indexer, LocalEmbeddingProvider, RAGEngine, Retriever

# ── Concrete Implementations ──────────────────────────────────────────────────

class OllamaLLMClient:
    """Ollama implementation of LLMClient protocol."""

    def __init__(self, host: str, model: str, options: dict):
        self.host = host
        self.model = model
        self.options = options
        import httpx
        self._client = httpx.AsyncClient(timeout=180.0)

    async def generate(self, prompt: str, system: str) -> str:
        try:
            response = await self._client.post(
                f"{self.host}/api/generate",
                json={
                    "model": self.model,
                    "prompt": prompt,
                    "system": system,
                    "stream": False,
                    "options": self.options,
                },
            )
            response.raise_for_status()
            return response.json().get("response", "Erreur : réponse Ollama vide")
        except Exception as e:
            return f"Erreur lors de la génération : {e}"

    async def check_health(self) -> bool:
        try:
            r = await self._client.get(f"{self.host}/api/tags")
            r.raise_for_status()
            return True
        except Exception:
            return False

    async def close(self):
        await self._client.aclose()


class PGVectorStore:
    """PostgreSQL + pgvector implementation of VectorStore protocol.

    Compose Retriever (recherche) + Indexer (indexation/stats/maintenance).
    Le pool DB est partagé via database.get_db() et distribué aux composants.
    """

    def __init__(self, db_url: str, embedding_provider: EmbeddingProvider):
        self.db_url = db_url
        self.embedding_provider = embedding_provider
        self._pool = None
        self._retriever = Retriever(embedding_provider)
        self._indexer = Indexer(embedding_provider)

    async def initialize(self) -> None:
        # Pool is created by database.get_db() at app startup
        self._pool = await get_db()
        self._retriever.attach_pool(self._pool)
        self._indexer.attach_pool(self._pool)

    async def retrieve(
        self,
        query_vector,
        top_k: int,
        threshold: float,
        metier_filter: str | None,
    ) -> list[dict]:
        if not self._pool:
            await self.initialize()
        return await self._retriever._search(
            query_vector, top_k, threshold, metier_filter
        )

    async def index_chunks(
        self,
        chunks: list[dict[str, object]],
        file_id: str,
        filename: str,
    ) -> None:
        if not self._pool:
            await self.initialize()
        await self._indexer.index_chunks(chunks, file_id, filename)

    async def get_stats(self) -> dict[str, object]:
        if not self._pool:
            await self.initialize()
        return await self._indexer.get_collection_stats()

    async def reset(self) -> None:
        if not self._pool:
            await self.initialize()
        await self._indexer.reset_collection()

    async def check_health(self) -> bool:
        if not self._pool:
            await self.initialize()
        try:
            await self._indexer.check_db_health()
            return True
        except Exception:
            return False

    async def close(self) -> None:
        await close_db()


class DefaultChunker:
    """RecursiveCharacterTextSplitter implementation of Chunker protocol."""

    def __init__(self, chunk_size: int = 400, chunk_overlap: int = 80):
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        self._splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", ". ", " ", ""],
            length_function=len,
        )

    def chunk(self, text: str, source: str, file_type: str) -> list[dict]:
        if not text.strip():
            return []
        raw_chunks = self._splitter.split_text(text)
        return [
            {
                "text": chunk,
                "metadata": {
                    "source": source,
                    "file_type": file_type,
                    "chunking_method": "recursive_char",
                    "chunk_size": len(chunk),
                },
            }
            for chunk in raw_chunks
            if chunk.strip()
        ]


class WhisperTranscriber:
    """Whisper implementation of Transcriber protocol (singleton)."""

    def __init__(self):
        self._model = None
        self._device = None

    def _get_model(self):
        import torch as _torch
        import whisper
        current_device = "cuda" if _torch.cuda.is_available() else "cpu"
        if self._model is None or self._device != current_device:
            if self._model is not None:
                import logging
                logging.getLogger(__name__).info("Whisper: rechargement (changement device)")
            import logging
            logging.getLogger(__name__).info(f"Chargement Whisper base ({current_device})...")
            self._model = whisper.load_model("base", device=current_device)
            self._device = current_device
            import logging
            logging.getLogger(__name__).info(f"Whisper base chargé sur {current_device}")
        return self._model

    def transcribe(self, file_path: str) -> str:
        try:
            model = self._get_model()
            fp16 = (self._device == "cuda")
            import logging
            logging.getLogger(__name__).info(f"Transcription {file_path} (fp16={fp16})...")
            result = model.transcribe(file_path, language="fr", fp16=fp16)
            text = result.get("text", "")
            import logging
            logging.getLogger(__name__).info(f"Transcription terminée ({len(text)} chars)")
            return text
        except ImportError:
            import logging
            logging.getLogger(__name__).warning("Whisper non installé — transcription ignorée")
            return "[Transcription non disponible — installer openai-whisper]"
        except Exception as e:
            import logging
            logging.getLogger(__name__).error(f"Transcription : {e}")
            return ""

    def unload(self) -> None:
        if self._model is not None:
            import torch as _torch
            del self._model
            self._model = None
            self._device = None
            if _torch.cuda.is_available():
                _torch.cuda.empty_cache()
            import logging
            logging.getLogger(__name__).info("Whisper déchargé (VRAM libérée)")


# ── Document Extractors (one per format) ──────────────────────────────────────

class PDFExtractor:
    @property
    def supported_extensions(self) -> set[str]:
        return {".pdf"}

    def extract(self, file_path: str) -> str:
        from pypdf import PdfReader
        text_parts = []
        try:
            with open(file_path, "rb") as f:
                reader = PdfReader(f)
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        text_parts.append(t)
        except Exception as e:
            import logging
            logging.getLogger(__name__).error(f"Extraction PDF : {e}")
        return "\n\n".join(text_parts)


class TextExtractor:
    @property
    def supported_extensions(self) -> set[str]:
        return {".txt", ".md"}

    def extract(self, file_path: str) -> str:
        try:
            with open(file_path, encoding="utf-8", errors="replace") as f:
                return f.read()
        except Exception as e:
            import logging
            logging.getLogger(__name__).error(f"Extraction texte : {e}")
            return ""


class DocxExtractor:
    @property
    def supported_extensions(self) -> set[str]:
        return {".docx"}

    def extract(self, file_path: str) -> str:
        from docx import Document as DocxDocument
        try:
            doc = DocxDocument(file_path)
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            import logging
            logging.getLogger(__name__).error(f"Extraction DOCX : {e}")
            return ""


class PptxExtractor:
    @property
    def supported_extensions(self) -> set[str]:
        return {".pptx"}

    def extract(self, file_path: str) -> str:
        from pptx import Presentation
        try:
            prs = Presentation(file_path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
            return "\n\n".join(parts)
        except Exception as e:
            import logging
            logging.getLogger(__name__).error(f"Extraction PPTX : {e}")
            return ""


class XlsxExtractor:
    @property
    def supported_extensions(self) -> set[str]:
        return {".xlsx"}

    def extract(self, file_path: str) -> str:
        import openpyxl
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
            rows = []
            for sheet in wb:
                for row in sheet.iter_rows(values_only=True):
                    r = " | ".join(str(c) for c in row if c is not None)
                    if r:
                        rows.append(r)
            wb.close()
            return "\n".join(rows)
        except Exception as e:
            import logging
            logging.getLogger(__name__).error(f"Extraction XLSX : {e}")
            return ""


class AudioVideoExtractor:
    """Uses WhisperTranscriber for audio/video."""

    def __init__(self, transcriber: WhisperTranscriber):
        self.transcriber = transcriber

    @property
    def supported_extensions(self) -> set[str]:
        return {".mp3", ".mp4", ".wav"}

    def extract(self, file_path: str) -> str:
        return self.transcriber.transcribe(file_path)


# ── Dependency Providers (FastAPI Depends) ────────────────────────────────────

@lru_cache
def get_settings_cached():
    return get_settings()


def get_embedding_provider() -> EmbeddingProvider:
    settings = get_settings_cached()
    return LocalEmbeddingProvider(settings.EMBEDDING_MODEL)


def get_llm_client() -> LLMClient:
    settings = get_settings_cached()
    options = {
        "temperature": settings.OLLAMA_TEMPERATURE,
        "top_p": settings.OLLAMA_TOP_P,
        "top_k": settings.OLLAMA_TOP_K,
        "num_predict": settings.OLLAMA_NUM_PREDICT,
        "num_ctx": settings.OLLAMA_NUM_CTX,
        "num_thread": settings.OLLAMA_NUM_THREAD,
        "num_gpu": settings.OLLAMA_NUM_GPU,
        "f16_kv": settings.OLLAMA_F16_KV,
    }
    return OllamaLLMClient(settings.OLLAMA_HOST, settings.OLLAMA_MODEL, options)


def get_rag_engine(
    embedding_provider: EmbeddingProvider = Depends(get_embedding_provider),
    llm_client: LLMClient = Depends(get_llm_client),
) -> RAGEngine:
    """Composition root : RAGEngine (facade) injecté via FastAPI Depends.

    Utilisé par main.py pour /chat (retrieve + generate), /documents, /indexing.
    """
    settings = get_settings_cached()
    return RAGEngine(
        db_url=settings.DATABASE_URL,
        embedding_provider=embedding_provider,
        llm_client=llm_client,
        ollama_host=settings.OLLAMA_HOST,
        model_name=settings.OLLAMA_MODEL,
    )


def get_vector_store(
    embedding_provider: EmbeddingProvider = Depends(get_embedding_provider),
) -> VectorStore:
    settings = get_settings_cached()
    return PGVectorStore(settings.DATABASE_URL, embedding_provider)


def get_chunker() -> Chunker:
    settings = get_settings_cached()
    return DefaultChunker(settings.CHUNK_SIZE, settings.CHUNK_OVERLAP)


def get_transcriber() -> Transcriber:
    return WhisperTranscriber()


def get_extractors(transcriber: Transcriber = Depends(get_transcriber)) -> list[DocumentExtractor]:
    return [
        PDFExtractor(),
        TextExtractor(),
        DocxExtractor(),
        PptxExtractor(),
        XlsxExtractor(),
        AudioVideoExtractor(transcriber),
    ]


def get_document_processor(
    chunker: Chunker = Depends(get_chunker),
    extractors: list[DocumentExtractor] = Depends(get_extractors),
    transcriber: Transcriber = Depends(get_transcriber),
) -> DocumentProcessor:
    settings = get_settings_cached()
    proc = DocumentProcessor.__new__(DocumentProcessor)
    proc.upload_dir = __import__("pathlib").Path(settings.UPLOAD_DIR)
    proc.upload_dir.mkdir(parents=True, exist_ok=True)
    proc._whisper_model = None
    proc._whisper_device = None
    # Override internal splitter with injected chunker
    proc._SPLITTER = chunker._splitter if hasattr(chunker, "_splitter") else None
    # Store injected dependencies
    proc._extractors = extractors
    proc._transcriber = transcriber
    proc._chunker = chunker
    return proc
